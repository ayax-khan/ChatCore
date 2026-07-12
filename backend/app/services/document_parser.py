import logging
import csv
import io
import re
from typing import Optional
from pathlib import Path

logger = logging.getLogger(__name__)


class DocumentParserService:
    def __init__(self):
        self.supported_extensions = ["pdf", "docx", "csv", "xlsx", "pptx", "txt", "md", "xml", "json"]

    def parse_file(self, file_path: str, mime_type: str | None = None) -> str:
        ext = file_path.rsplit(".", 1)[-1].lower() if "." in file_path else ""
        if ext not in self.supported_extensions:
            return f"Unsupported file type: {ext}"
        try:
            if ext == "pdf":
                return self._parse_pdf(file_path)
            elif ext == "docx":
                return self._parse_docx(file_path)
            elif ext == "csv":
                return self._parse_csv(file_path)
            elif ext == "xlsx":
                return self._parse_xlsx(file_path)
            elif ext == "pptx":
                return self._parse_pptx(file_path)
            elif ext in ("txt", "md"):
                return self._parse_text(file_path)
            elif ext in ("xml", "json"):
                return self._parse_code_file(file_path)
            else:
                return self._parse_text(file_path)
        except Exception as e:
            logger.error(f"Failed to parse file {file_path}: {e}")
            return ""

    def parse_bytes(self, data: bytes, ext: str, custom_mime: str | None = None) -> str:
        import tempfile
        if ext not in self.supported_extensions:
            return ""
        try:
            if ext == "pdf":
                return self._parse_pdf_bytes(data)
            elif ext == "docx":
                return self._parse_docx_bytes(data)
            elif ext == "csv":
                return self._parse_csv_bytes(data)
            elif ext == "xlsx":
                return self._parse_xlsx_bytes(data)
            elif ext == "pptx":
                return self._parse_pptx_bytes(data)
            elif ext in ("txt", "md"):
                return data.decode("utf-8", errors="replace")
            elif ext in ("xml", "json"):
                return data.decode("utf-8", errors="replace")
            else:
                return data.decode("utf-8", errors="replace")
        except Exception as e:
            logger.error(f"Failed to parse data as {ext}: {e}")
            return ""

    def _parse_pdf(self, file_path: str) -> str:
        try:
            import PyPDF2
        except ImportError:
            return self._parse_with_pdfminer(file_path)
        try:
            text = ""
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for i, page in enumerate(reader.pages):
                    page_text = page.extract_text() or ""
                    if page_text.strip():
                        text += f"\n\n--- Page {i+1} ---\n\n{page_text}"
            return text.strip()
        except Exception as e:
            return self._parse_with_pdfminer(file_path)

    def _parse_with_pdfminer(self, file_path: str) -> str:
        try:
            from pdfminer.high_level import extract_text
            text = extract_text(file_path)
            return text.strip()
        except ImportError:
            return "[PDF parsing not available: install PyPDF2 or pdfminer.six]"

    def _parse_pdf_bytes(self, data: bytes) -> str:
        try:
            import PyPDF2
            text = ""
            reader = PyPDF2.PdfReader(data)
            for i, page in enumerate(reader.pages):
                page_text = page.extract_text() or ""
                if page_text.strip():
                    text += f"\n\n--- Page {i+1} ---\n\n{page_text}"
            return text.strip()
        except ImportError:
            return "[PDF parsing not available: install PyPDF2]"

    def _parse_docx(self, file_path: str) -> str:
        try:
            from docx import Document
        except ImportError:
            return "[DOCX parsing not available: install python-docx]"
        try:
            doc = Document(file_path)
            text_parts = []
            for para in doc.paragraphs:
                text_parts.append(para.text)
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text for cell in row.cells)
                    text_parts.append(row_text)
            return "\n".join(text_parts)
        except Exception as e:
            logger.error(f"Failed to parse DOCX {file_path}: {e}")
            return ""

    def _parse_docx_bytes(self, data: bytes) -> str:
        try:
            from docx import Document
        except ImportError:
            return "[DOCX parsing not available: install python-docx]"
        try:
            import io
            doc = Document(io.BytesIO(data))
            return "\n".join(p.text for p in doc.paragraphs)
        except Exception as e:
            logger.error(f"Failed to parse DOCX data: {e}")
            return ""

    def _parse_csv(self, file_path: str) -> str:
        try:
            rows = []
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                reader = csv.reader(f)
                for row in reader:
                    rows.append(" | ".join(row))
            return "\n".join(rows)
        except Exception as e:
            logger.error(f"Failed to parse CSV {file_path}: {e}")
            return ""

    def _parse_csv_bytes(self, data: bytes) -> str:
        try:
            rows = []
            text = data.decode("utf-8", errors="replace")
            reader = csv.reader(io.StringIO(text))
            for row in reader:
                rows.append(" | ".join(row))
            return "\n".join(rows)
        except Exception as e:
            logger.error(f"Failed to parse CSV data: {e}")
            return ""

    def _parse_xlsx(self, file_path: str) -> str:
        try:
            import openpyxl
        except ImportError:
            return "[XLSX parsing not available: install openpyxl]"
        try:
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            text_parts = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                text_parts.append(f"\n--- Sheet: {sheet_name} ---\n")
                for row in ws.iter_row(values_only=True):
                    row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                    if row_text.strip():
                        text_parts.append(row_text)
            return "\n".join(text_parts)
        except Exception as e:
            logger.error(f"Failed to parse XLSX {file_path}: {e}")
            return ""

    def _parse_xlsx_bytes(self, data: bytes) -> str:
        try:
            import openpyxl
        except ImportError:
            return "[XLSX parsing not available: install openpyxl]"
        try:
            wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
            text_parts = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                text_parts.append(f"\n--- Sheet: {sheet_name} ---\n")
                for row in ws.iter_row(values_only=True):
                    row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                    if row_text.strip():
                        text_parts.append(row_text)
            return "\n".join(text_parts)
        except Exception as e:
            logger.error(f"Failed to parse XLSX data: {e}")
            return ""

    def _parse_pptx(self, file_path: str) -> str:
        try:
            from pptx import Presentation
        except ImportError:
            return "[PPTX parsing not available: install python-pptx]"
        try:
            prs = Presentation(file_path)
            text_parts = []
            for i, slide in enumerate(prs.slides):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text)
                if slide_texts:
                    text_parts.append(f"\n--- Slide {i+1} ---\n" + "\n".join(slide_texts))
            return "\n".join(text_parts)
        except Exception as e:
            logger.error(f"Failed to parse PPTX {file_path}: {e}")
            return ""

    def _parse_pptx_bytes(self, data: bytes) -> str:
        try:
            from pptx import Presentation
        except ImportError:
            return "[PPTX parsing not available: install python-pptx]"
        try:
            prs = Presentation(io.BytesIO(data))
            text_parts = []
            for i, slide in enumerate(prs.slides):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text)
                if slide_texts:
                    text_parts.append(f"\n--- Slide {i+1} ---\n" + "\n".join(slide_texts))
            return "\n".join(text_parts)
        except Exception as e:
            logger.error(f"Failed to parse PPTX data: {e}")
            return ""

    def _parse_text(self, file_path: str) -> str:
        try:
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                return f.read()
        except Exception as e:
            logger.error(f"Failed to read text file {file_path}: {e}")
            return ""

    def _parse_code_file(self, file_path: str) -> str:
        try:
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                content = f.read()
            ext = file_path.rsplit(".", 1)[-1]
            return f"```{ext}\n{content}\n```"
        except Exception as e:
            logger.error(f"Failed to read code file {file_path}: {e}")
            return ""

    def detect_language(self, text: str) -> str:
        if not text or len(text) < 20:
            return "en"
        char_count = len(text)
        chinese = len(re.findall(r"[\u4e00-\u9fff]", text))
        japanese = len(re.findall(r"[\u3040-\u309f\u30a0-\u30ff]", text))
        korean = len(re.findall(r"[\uac00-\ud7af]", text))
        arabic = len(re.findall(r"[\u0600-\u06ff]", text))
        devanagari = len(re.findall(r"[\u0900-\u097f]", text))
        cyrillic = len(re.findall(r"[\u0400-\u04ff]", text))
        if chinese / char_count > 0.3:
            if japanese > chinese: return "ja"
            if korean > chinese: return "ko"
            return "zh"
        if korean / char_count > 0.3: return "ko"
        if arabic / char_count > 0.3: return "ar"
        if devanagari / char_count > 0.3: return "hi"
        if cyrillic / char_count > 0.3: return "ru"
        if japanese / char_count > 0.1: return "ja"
        return "en"
